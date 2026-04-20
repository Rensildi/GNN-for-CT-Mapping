"""Generate the project status PowerPoint deck.

Run:
    source AI/bin/activate
    python presentation/build_presentation.py

Output:
    presentation/gnn_lung_nodule_status.pptx

The deck is 16:9, light-themed (white background, deep-blue titles, dark-
gray body), Calibri throughout. All slides use the blank layout so
styling stays in our control — no theme fight with PowerPoint defaults.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt


# --- Palette (light theme) -------------------------------------------------
DEEP_BLUE = RGBColor(0x1F, 0x4E, 0x79)   # titles
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6) # dividers / accents
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)   # body text
MID_GRAY = RGBColor(0x6B, 0x6B, 0x6B)    # subtitles / captions
PLACEHOLDER_GRAY = RGBColor(0xB0, 0xB0, 0xB0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# --- Layout grid -----------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

TITLE_LEFT = Inches(0.5)
TITLE_TOP = Inches(0.35)
TITLE_W = Inches(12.3)
TITLE_H = Inches(0.9)

DIVIDER_TOP = Inches(1.25)

BODY_LEFT = Inches(0.6)
BODY_TOP = Inches(1.5)
BODY_W = Inches(12.1)
BODY_H = Inches(5.7)


REPO_ROOT = Path(__file__).resolve().parent.parent
ARCHITECTURE_PNG = REPO_ROOT / "GNN_for_CT_Mapping/experiments/harrison/architecture_exp1.png"


# --- Helpers ---------------------------------------------------------------

def _set_run(run, size_pt: float, *, bold: bool = False, color: RGBColor = DARK_GRAY,
             font_name: str = "Calibri") -> None:
    """Apply a consistent font style to a run."""
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color


def _new_slide(prs: Presentation):
    """Blank slide with a solid white background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    # Force white background even if the theme has something else.
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return slide


def _add_title(slide, text: str, *, accent_divider: bool = True) -> None:
    box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_W, TITLE_H)
    tf = box.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        _set_run(run, 30, bold=True, color=DEEP_BLUE)
    if accent_divider:
        # Thin accent line under the title to give slides a consistent skeleton.
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, TITLE_LEFT, DIVIDER_TOP, Inches(1.6), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_BLUE
        line.line.fill.background()


def _add_bulleted_body(slide, items: list[tuple[int, str]], *, top: Inches = BODY_TOP,
                       height: Inches = BODY_H, body_font_size: float = 18) -> None:
    """Add a bullet list where each item is (indent_level, text).

    indent_level 0 renders with a round bullet, 1 with an en-dash, 2 with
    a smaller dot. We draw bullets manually (rather than relying on the
    default template's list styles) so the look stays consistent across
    builds and PowerPoint versions.
    """
    box = slide.shapes.add_textbox(BODY_LEFT, top, BODY_W, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    bullet_chars = {0: "•", 1: "–", 2: "·"}
    # marL (text left) and indent (bullet offset left of marL) per nesting
    # level, in EMU. Applying these via OXML gives proper hanging indent
    # so wrapped continuation lines align with the text of the first line,
    # not with the bullet character.
    mar_l_emu = {0: Inches(0.35), 1: Inches(0.85), 2: Inches(1.35)}
    hanging_emu = Inches(0.3)

    for idx, (level, text) in enumerate(items):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        marker = bullet_chars.get(level, "•")
        p.text = f"{marker}  {text}"
        p.level = level
        # Attach paragraph margin + hanging indent directly in OXML —
        # python-pptx's public API doesn't expose marL / indent yet.
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(int(mar_l_emu.get(level, Inches(0.35)))))
        pPr.set("indent", str(-int(hanging_emu)))
        for run in p.runs:
            _set_run(run, body_font_size - level * 2, color=DARK_GRAY)


def _add_centered_paragraph(slide, text: str, *, top_inches: float, height_inches: float,
                             size: float, color: RGBColor = DARK_GRAY, bold: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top_inches), Inches(12.3), Inches(height_inches))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        _set_run(run, size, bold=bold, color=color)


def _add_placeholder_box(slide, top_inches: float, height_inches: float, caption: str) -> None:
    """Draw a dashed placeholder rectangle with a caption."""
    left = Inches(1.5)
    width = Inches(10.33)
    top = Inches(top_inches)
    height = Inches(height_inches)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    shape.line.color.rgb = PLACEHOLDER_GRAY
    shape.line.width = Pt(1.5)
    # Caption centered inside the box.
    tf = shape.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = caption
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        _set_run(run, 18, color=MID_GRAY, bold=False)


def _add_footer(slide, idx: int, total: int) -> None:
    """Small page marker bottom-right — helps during a live presentation."""
    box = slide.shapes.add_textbox(Inches(11.5), Inches(7.1), Inches(1.6), Inches(0.3))
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = f"{idx} / {total}"
    p.alignment = PP_ALIGN.RIGHT
    for run in p.runs:
        _set_run(run, 10, color=MID_GRAY)


# --- Slide builders --------------------------------------------------------

def build_title_slide(prs: Presentation) -> None:
    slide = _new_slide(prs)
    # Project title — sized to fit the 13.3-inch slide width on one line so
    # it never collides with the subtitle below.
    _add_centered_paragraph(slide,
        "GNN-Based Lung Nodule Malignancy Classification",
        top_inches=2.35, height_inches=0.9, size=34, color=DEEP_BLUE, bold=True)
    _add_centered_paragraph(slide,
        "Multi-Modal Feature Embedding on CT Scans",
        top_inches=3.35, height_inches=0.7, size=24, color=DARK_GRAY)
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.67), Inches(4.25), Inches(2.0), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()
    # Authors + course
    _add_centered_paragraph(slide,
        "Swathi Gopal    •    Harrison Lavins    •    Rensildi Kalanxhi",
        top_inches=4.55, height_inches=0.6, size=20, color=DARK_GRAY)
    _add_centered_paragraph(slide,
        "CSC 7760 — Deep Learning    |    Status Update — April 2026",
        top_inches=5.25, height_inches=0.6, size=16, color=MID_GRAY)


def build_project_goals(prs):
    slide = _new_slide(prs)
    _add_title(slide, "Project Goals")
    _add_bulleted_body(slide, [
        (0, "Primary question: Does modeling inter-nodule relationships as a similarity graph"),
        (1, "improve malignancy classification versus classifying each nodule independently?"),
        (0, "Head-to-head comparison on LIDC-IDRI"),
        (1, "2-layer Graph Convolutional Network head"),
        (1, "Parameter-matched MLP baseline on identical multi-modal features"),
        (0, "Cross-dataset generalization test: LIDC-IDRI → LUNA25"),
        (0, "Validation against pathology-confirmed subset (~157 nodules with biopsy / surgery ground truth)"),
        (0, "Keep architecture within an RTX 3060 (12 GB VRAM) training budget"),
    ])


def build_background(prs):
    slide = _new_slide(prs)
    _add_title(slide, "Background")
    _add_bulleted_body(slide, [
        (0, "Lung cancer is the leading cause of cancer mortality worldwide"),
        (0, "CT-based nodule classifiers reach AUC 0.88 – 0.96 on LIDC-IDRI"),
        (1, "But they classify each nodule independently, ignoring inter-case similarity"),
        (0, "Graph Neural Networks propagate information across similar cases via message passing"),
        (1, "Patient-similarity GNNs help in cancer subtyping and prognosis"),
        (1, "Not yet applied to LIDC lung nodule malignancy classification on CT — this project fills that gap"),
        (0, "Closest prior work: Ma et al. 2023 (AUC 0.9629 on LIDC)"),
        (1, "Uses GCN for cross-CNN feature fusion, not for node classification on a similarity graph"),
    ])


def build_proposed_approach(prs):
    slide = _new_slide(prs)
    _add_title(slide, "Proposed Approach — Two-Stage Pipeline")
    _add_bulleted_body(slide, [
        (0, "Stage 1 — frozen multi-modal feature extraction"),
        (1, "Image: 48³-voxel CT patch → frozen Med3D ResNet-50 → 2048-D pooled features → Linear → 256-D"),
        (1, "Clinical: 8 LIDC-IDRI attributes → per-attribute nn.Embedding tables → 64-D"),
        (1, "Spatial: nodule centroid (x, y, z) in mm → sinusoidal positional encoding → 48-D"),
        (1, "Concatenate → LayerNorm → Linear → 256-D unified node feature"),
        (0, "Stage 2 — trained node-classification head"),
        (1, "GCN: 2 × GCNConv layers over the cohort-wide KNN graph, with dropout + Linear head"),
        (1, "MLP baseline: same widths and dropout, no graph structure — isolates the graph effect"),
        (0, "Training: weighted cross-entropy, Adam (lr 1e-3), 5-fold patient-level CV, early stop on val AUC"),
    ])


def build_architecture_diagram(prs):
    slide = _new_slide(prs)
    _add_title(slide, "Architecture")
    # Place the image centered below the title.
    if ARCHITECTURE_PNG.exists():
        # Target height ~ 5.6 in keeps it under the slide bottom; width scales.
        pic = slide.shapes.add_picture(str(ARCHITECTURE_PNG), Inches(0.0), Inches(1.45), height=Inches(5.8))
        # Center horizontally after PowerPoint computes the width.
        pic.left = int((SLIDE_W - pic.width) / 2)
    else:
        _add_placeholder_box(slide, 2.0, 4.5,
                             f"(missing architecture_exp1.png — regenerate with draw_architecture.py)")


def build_datasets(prs):
    slide = _new_slide(prs)
    _add_title(slide, "Datasets")
    _add_bulleted_body(slide, [
        (0, "LIDC-IDRI (primary — training + within-dataset eval)"),
        (1, "1,010 patients  •  1,308 CT studies  •  ~133 GB DICOM"),
        (1, "≈ 7,371 nodules ≥ 3 mm, up to 4 radiologist readings per nodule"),
        (1, "8 descriptive attributes + 1-5 malignancy rating per reader per nodule"),
        (1, "~157 pathology-confirmed diagnoses (biopsy / surgery / follow-up)"),
        (1, "After preprocessing + binarization: 1,128 nodules across 588 patients (805 benign / 323 malignant)"),
        (0, "LUNA25 (secondary — cross-dataset generalization, Experiment 4)"),
        (1, "4,096 CT exams  •  6,163 nodules from National Lung Screening Trial"),
        (1, "Binary malignancy labels from clinical follow-up"),
        (1, "Released for MICCAI 2025 (Zenodo, CC BY 4.0)"),
    ])


def build_dataset_viz_placeholder(prs, title: str, caption: str):
    slide = _new_slide(prs)
    _add_title(slide, title)
    _add_placeholder_box(slide, top_inches=1.6, height_inches=5.3, caption=caption)


def build_preprocessing(prs):
    slide = _new_slide(prs)
    _add_title(slide, "Preprocessing & Data Cleaning")
    _add_bulleted_body(slide, [
        (0, "XML annotations — stdlib ElementTree parser (no pylidc)"),
        (1, "Malignancy rating tracked in a separate field from the 8 feature attributes to prevent label leakage"),
        (0, "DICOM series loading via pydicom"),
        (1, "Slices ordered by physical Z (ImagePositionPatient[2]) — robust to vendor-reversed InstanceNumber"),
        (1, "Per-slice RescaleSlope / RescaleIntercept applied to get Hounsfield Units"),
        (0, "Cross-reader nodule clustering"),
        (1, "Greedy centroid-distance match (8 mm threshold) across up to 4 radiologists per scan"),
        (1, "Follow-up: 6 / 1,128 clusters have >4 merged readers → retighten or switch to IoU-based clustering"),
        (0, "Patch extraction & normalization (feature-time)"),
        (1, "48³-voxel patches centered on each nodule, HU-clipped to [-1000, 400] and normalized to [0, 1]"),
        (1, "Resampled to isotropic 1 mm³ before Med3D inference"),
        (0, "Label binarization: mean malignancy ≤ 2 → benign (0), ≥ 4 → malignant (1), ~3 excluded"),
        (0, "Patient-level 5-fold StratifiedGroupKFold — no patient straddles train / val"),
    ])


def build_novelty(prs):
    slide = _new_slide(prs)
    _add_title(slide, "What's Novel")
    _add_bulleted_body(slide, [
        (0, "First node-classification GCN applied to LIDC malignancy"),
        (1, "Prior GCN work on LIDC uses the GCN for cross-CNN feature fusion, not for inter-nodule reasoning"),
        (0, "Multi-modal node features fuse three complementary sources"),
        (1, "Imaging (Med3D), clinical (8 radiologist attributes), spatial (sinusoidal positional encoding)"),
        (0, "Inductive evaluation protocol for the cohort-wide graph"),
        (1, "Edges fit on train-only nodes; val nodes inserted at eval time with edges into train neighbors only"),
        (1, "No val-to-val edges — eliminates the standard transductive-leak failure mode under patient-level CV"),
        (0, "Parameter-matched MLP control isolates the effect of message passing from raw capacity"),
        (0, "Committed encoder follow-up: FMCIB (Aerts lab, 2024) as a second encoder axis in Experiment 3"),
    ])


def build_experiments_overview(prs):
    slide = _new_slide(prs)
    _add_title(slide, "Experiments")
    _add_bulleted_body(slide, [
        (0, "Experiment 1 — GCN vs. MLP baseline on LIDC-IDRI"),
        (1, "Status: pipeline complete, training blocked on Med3D checkpoint download"),
        (0, "Experiment 2 — Graph construction ablation"),
        (1, "8-cell grid: k ∈ {5, 10, 15, 20} × metric ∈ {cosine, euclidean}"),
        (1, "Best (k, metric) adopted as default for downstream experiments"),
        (0, "Experiment 3 — Feature-modality × encoder ablation"),
        (1, "6-cell 3×2 grid: {image, image+attrs, image+attrs+spatial} × {Med3D, FMCIB}"),
        (1, "Pathology-confirmed subset is primary validation axis"),
        (0, "Experiment 4 — Cross-dataset generalization to LUNA25"),
        (1, "Apply LIDC-trained model to LUNA25 nodules (different scanners, protocols, patients)"),
    ])


def build_experiment1(prs):
    slide = _new_slide(prs)
    _add_title(slide, "Experiment 1 — GCN vs. MLP Baseline")
    _add_bulleted_body(slide, [
        (0, "Objective: head-to-head AUC test of graph-based aggregation versus independent classification"),
        (0, "Hypothesis (H1): GCN macro-AUC > MLP macro-AUC by ≥ 0.01 (paired Wilcoxon p < 0.05, 5 folds)"),
        (0, "Both models consume identical 256-D fused features — the only difference is the KNN graph"),
        (0, "Inductive KNN construction prevents train / val leakage"),
        (0, "Metrics: AUC, AUPRC, sensitivity / specificity at Youden-J, Brier score, 10-bin reliability diagram"),
        (0, "Per-fold class balance: val malignant counts 67, 77, 63, 62, 54 (all above the plan's ≥ 30 threshold)"),
        (0, "Training: Adam (lr 1e-3, wd 1e-4), 100 epochs with patience-15 early stop on val AUC"),
        (0, "Status: code + splits committed; waiting on Tencent/MedicalNet resnet_50.pth download"),
    ])


def build_experiment2(prs):
    slide = _new_slide(prs)
    _add_title(slide, "Experiment 2 — Graph Construction Ablation  (Planned)")
    _add_bulleted_body(slide, [
        (0, "8-cell grid: k ∈ {5, 10, 15, 20}  ×  metric ∈ {cosine, euclidean}"),
        (0, "Fixed GCN architecture from Experiment 1; re-fit KNN per cell per fold"),
        (0, "Decision rule: Bonferroni-corrected paired Wilcoxon across the 8 cells"),
        (0, "Best cell becomes the default for Experiments 3 and 4; null result → keep (k=10, cosine)"),
    ])
    _add_placeholder_box(slide, top_inches=5.3, height_inches=1.8,
                         caption="Results to follow — AUC heatmap over (k, metric)")


def build_experiment3(prs):
    slide = _new_slide(prs)
    _add_title(slide, "Experiment 3 — Feature × Encoder Ablation  (Planned)")
    _add_bulleted_body(slide, [
        (0, "3 × 2 grid: feature configs × image encoders"),
        (1, "Features: {image only, image + attributes, image + attributes + spatial}"),
        (1, "Encoders:  {Med3D ResNet-50 (baseline), FMCIB foundation model (follow-up)}"),
        (0, "Pathology-confirmed subset is the primary validation axis (true ground truth)"),
        (0, "Tests whether radiologist attributes add unique signal beyond the image branch (leakage check)"),
    ])
    _add_placeholder_box(slide, top_inches=5.3, height_inches=1.8,
                         caption="Results to follow — AUC heatmap over feature-config × encoder")


def build_experiment4(prs):
    slide = _new_slide(prs)
    _add_title(slide, "Experiment 4 — LUNA25 Cross-Dataset  (Planned)")
    _add_bulleted_body(slide, [
        (0, "Apply LIDC-trained model to LUNA25 nodules (held-out entirely from training)"),
        (0, "Different source patients, scanners, acquisition protocols"),
        (0, "Label source: clinical follow-up (binary), not radiologist consensus — cleaner ground truth"),
        (0, "Tests whether learned representations generalize beyond LIDC's specific reader distribution"),
    ])
    _add_placeholder_box(slide, top_inches=5.3, height_inches=1.8,
                         caption="Results to follow — LIDC→LUNA25 transfer performance")


def build_teammate_placeholder(prs, name: str):
    slide = _new_slide(prs)
    _add_title(slide, f"{name}'s Work")
    _add_placeholder_box(slide, top_inches=1.8, height_inches=5.2,
                         caption=f"(Placeholder — {name} to fill in focus area, methods, and current progress)")


def build_current_results(prs):
    slide = _new_slide(prs)
    _add_title(slide, "Current Results")
    _add_bulleted_body(slide, [
        (0, "Data pipeline complete"),
        (1, "1,128 labeled nodules from 588 patients (805 benign / 323 malignant)"),
        (1, "Patient-level 5-fold CV splits committed; no train / val patient leakage"),
        (0, "Model code complete and smoke-tested"),
        (1, "Multi-modal fusion, 2-layer GCN, parameter-matched MLP, inductive KNN construction"),
        (1, "Med3D ResNet-50 backbone wired; forward pass verified at (N, 2048) on CUDA (random init)"),
        (0, "No training numbers yet"),
        (1, "Training requires Tencent/MedicalNet pretrained weights (manual download — no wget URL)"),
        (1, "Full GCN vs. MLP run expected to take ~1 – 2 hours on the RTX 3060 once features are cached"),
    ])


def build_next_steps(prs):
    slide = _new_slide(prs)
    _add_title(slide, "Next Steps")
    _add_bulleted_body(slide, [
        (0, "Download Tencent/MedicalNet resnet_50.pth"),
        (0, "Verify checkpoint: python -m GNN_for_CT_Mapping.experiments.harrison.scripts.verify_med3d ..."),
        (0, "Cache Med3D features across 1,128 nodules (~30 min on the 3060)"),
        (0, "Run train_exp1.py across all 5 folds; monitor AUC / loss curves in TensorBoard"),
        (0, "Collect paired-Wilcoxon GCN vs. MLP test on AUC / AUPRC; write results_exp1.md"),
        (0, "Investigate data-quality follow-ups"),
        (1, "6 over-merged nodule clusters at the 8 mm centroid threshold"),
        (1, "Mean-malignancy edge case (zero-reader nodules in parquet)"),
        (0, "Begin Experiment 2 (graph construction ablation) — reuses Experiment 1's cached features"),
        (0, "Coordinate teammate deliverables for Rensildi and Swathi"),
    ])


# --- Top-level orchestration ----------------------------------------------

def build(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slides appended in order; the footer numbers get patched in after the
    # full deck is built (we need the final total count first).
    builders = [
        build_title_slide,
        build_project_goals,
        build_background,
        build_proposed_approach,
        build_architecture_diagram,
        build_datasets,
        lambda p: build_dataset_viz_placeholder(p,
            "Dataset Visualization — Nodule Size Distribution",
            "Histogram of nodule diameters across LIDC-IDRI (to be added after preprocess run)"),
        lambda p: build_dataset_viz_placeholder(p,
            "Dataset Visualization — Example 48³ CT Patches",
            "Grid of representative benign and malignant patches (axial mid-slice)"),
        lambda p: build_dataset_viz_placeholder(p,
            "Dataset Visualization — Attribute Correlations",
            "Heatmap of the 8 LIDC attribute means vs. malignancy label"),
        build_preprocessing,
        build_novelty,
        build_experiments_overview,
        build_experiment1,
        build_experiment2,
        build_experiment3,
        build_experiment4,
        lambda p: build_teammate_placeholder(p, "Rensildi"),
        lambda p: build_teammate_placeholder(p, "Swathi"),
        build_current_results,
        build_next_steps,
    ]
    for builder in builders:
        builder(prs)

    # Add per-slide footers now that we know the final count.
    total = len(prs.slides)
    for idx, slide in enumerate(prs.slides, start=1):
        _add_footer(slide, idx, total)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"Wrote {out_path}  ({total} slides)")


if __name__ == "__main__":
    import sys
    default = Path(__file__).resolve().parent / "gnn_lung_nodule_status.pptx"
    target = Path(sys.argv[1]) if len(sys.argv) > 1 else default
    build(target)
